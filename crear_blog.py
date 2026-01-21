# ============================================================
# crear_blog.py — Crear Blog (SOLO CREACIÓN) — MEJORADO
# - Mejor extracción de imágenes (PDF/DOCX/PPTX) con filtros
# - Detección de tablas desde DOCX/PPTX/PDF (pdfplumber opcional)
# - UI para previsualizar e insertar tablas detectadas en el documento
# ============================================================

import streamlit as st
from markitdown import MarkItDown

import os
import re
import io
import json
import shutil
import zipfile
import tempfile
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

import pandas as pd
from streamlit_ace import st_ace

# PDF viewer (opcional)
try:
    from streamlit_pdf_viewer import pdf_viewer
except Exception:
    pdf_viewer = None


BLOG_DIR = Path("blog_posts")
INDEX_FILE = BLOG_DIR / "index.json"
BLOG_DIR.mkdir(exist_ok=True)

CURSOR_TOKEN = "<<CURSOR>>"


# =============================
# INDEX helpers
# =============================
def load_index():
    if not INDEX_FILE.exists():
        return []
    try:
        return json.loads(INDEX_FILE.read_text(encoding="utf-8"))
    except Exception:
        return []


def save_index(items):
    INDEX_FILE.write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")


def upsert_post_meta(meta: dict):
    items = load_index()
    items = [x for x in items if x.get("id") != meta.get("id")]
    items.insert(0, meta)
    save_index(items)


def zip_folder(folder: Path) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(folder):
            for name in files:
                p = Path(root) / name
                arc = p.relative_to(folder).as_posix()
                z.write(p, arcname=arc)
    buf.seek(0)
    return buf.getvalue()


def slugify(text: str) -> str:
    text = (text or "").strip().lower()
    text = re.sub(r"[^\w\s-]", "", text, flags=re.UNICODE)
    text = re.sub(r"[\s_-]+", "-", text).strip("-")
    return text or "post"


# =============================
# MARKDOWN helpers
# =============================
def insert_at_cursor_token(md_text: str, insert_text: str) -> str:
    md = md_text or ""
    if CURSOR_TOKEN in md:
        return md.replace(CURSOR_TOKEN, insert_text + CURSOR_TOKEN, 1)
    return md + "\n\n" + insert_text


def ensure_cursor_token(md_text: str) -> str:
    md = md_text or ""
    if CURSOR_TOKEN in md:
        return md
    return md + "\n\n" + CURSOR_TOKEN + "\n"


def remove_all_cursor_tokens(md_text: str) -> str:
    return (md_text or "").replace(CURSOR_TOKEN, "")


def make_image_snippet(img_name: str):
    return f"""![{img_name}](images/{img_name})

*Figura X. Descripción de la figura*

"""


def infer_title_from_md(md_text: str, fallback: str) -> str:
    for line in (md_text or "").splitlines():
        if line.strip().startswith("# "):
            t = line.strip()[2:].strip()
            return t if t else fallback
    return fallback


def infer_description_from_md(md_text: str, max_len: int = 220) -> str:
    lines = [l.rstrip() for l in (md_text or "").splitlines()]
    buff = []
    in_code = False

    for l in lines:
        if l.strip().startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue

        if not l.strip():
            if buff:
                para = " ".join(buff).strip()
                if para:
                    return (para[: max_len - 1] + "…") if len(para) > max_len else para
                buff = []
            continue

        s = l.strip()
        if s.startswith("#") or s.startswith(">"):
            continue
        if s.startswith("- ") or s.startswith("* ") or re.match(r"^\d+\.", s):
            continue
        buff.append(s)

    para = " ".join(buff).strip()
    if not para:
        return ""
    return (para[: max_len - 1] + "…") if len(para) > max_len else para


# =============================
# Table helpers
# =============================
def df_to_markdown_table(df: pd.DataFrame, max_rows: int = 200) -> str:
    if df is None or df.empty:
        return ""
    df2 = df.copy()
    if len(df2) > max_rows:
        df2 = df2.head(max_rows)
    # evitar columnas "Unnamed"
    df2.columns = [("" if str(c).startswith("Unnamed") else str(c)) for c in df2.columns]
    return df2.to_markdown(index=False)


def extract_table_from_excel(file_path: str, sheet_name: str | None):
    return pd.read_excel(file_path, sheet_name=sheet_name, engine="openpyxl")


def extract_table_from_csv(file_path: str):
    return pd.read_csv(file_path)


def make_table_snippet(df: pd.DataFrame, title: str = "Tabla"):
    md_table = df_to_markdown_table(df)
    if not md_table.strip():
        return ""
    return f"""**{title}.** *(completa/edita el título si quieres)*

{md_table}

"""


# =============================
# NUEVO: extracción de TABLAS desde documentos
# =============================
def _clean_table_rows(rows: List[List[Any]]) -> List[List[str]]:
    # normaliza valores a string, recorta espacios
    cleaned = []
    for r in rows:
        rr = []
        for x in r:
            s = "" if x is None else str(x)
            rr.append(s.strip())
        # descarta filas completamente vacías
        if any(v != "" for v in rr):
            cleaned.append(rr)
    return cleaned


def _rows_to_df(rows: List[List[str]]) -> Optional[pd.DataFrame]:
    if not rows or len(rows) < 2:
        return None

    # iguala largo de filas
    max_len = max(len(r) for r in rows)
    norm = [r + [""] * (max_len - len(r)) for r in rows]

    header = norm[0]
    # si header está casi vacío, usa nombres genéricos
    if sum(1 for h in header if h.strip() != "") <= max(1, max_len // 4):
        cols = [f"col_{i+1}" for i in range(max_len)]
        data = norm
    else:
        cols = [h if h.strip() else f"col_{i+1}" for i, h in enumerate(header)]
        data = norm[1:]

    df = pd.DataFrame(data, columns=cols)
    # elimina filas totalmente vacías
    df = df.loc[~(df.apply(lambda r: all(str(x).strip() == "" for x in r), axis=1))]
    if df.empty:
        return None
    return df


def extract_tables_docx(path: str, max_tables: int = 50) -> List[Dict[str, Any]]:
    try:
        from docx import Document
    except Exception:
        return []

    out = []
    try:
        doc = Document(path)
        for ti, table in enumerate(doc.tables[:max_tables], start=1):
            rows = []
            for row in table.rows:
                rows.append([cell.text for cell in row.cells])
            rows = _clean_table_rows(rows)
            df = _rows_to_df(rows)
            if df is None or df.empty:
                continue
            out.append({"name": f"Tabla_DOCX_{ti:02d}", "df": df})
    except Exception:
        return []
    return out


def extract_tables_pptx(path: str, max_tables: int = 50) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except Exception:
        return []

    out = []
    try:
        prs = Presentation(path)
        tcount = 0
        for si, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                try:
                    if not hasattr(shape, "has_table") or not shape.has_table:
                        continue
                    tbl = shape.table
                    rows = []
                    for r in range(len(tbl.rows)):
                        row = []
                        for c in range(len(tbl.columns)):
                            row.append(tbl.cell(r, c).text)
                        rows.append(row)
                    rows = _clean_table_rows(rows)
                    df = _rows_to_df(rows)
                    if df is None or df.empty:
                        continue
                    tcount += 1
                    out.append({"name": f"Tabla_PPTX_S{si:02d}_{tcount:02d}", "df": df})
                    if tcount >= max_tables:
                        return out
                except Exception:
                    continue
    except Exception:
        return []
    return out


def extract_tables_pdf(path: str, max_tables: int = 80) -> List[Dict[str, Any]]:
    """
    Requiere pdfplumber para tablas más confiables.
    Si no está instalado, retorna [] sin romper.
    """
    try:
        import pdfplumber
    except Exception:
        return []

    out = []
    try:
        with pdfplumber.open(path) as pdf:
            tcount = 0
            for pi, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables()
                except Exception:
                    tables = []
                for tj, t in enumerate(tables, start=1):
                    if not t or len(t) < 2:
                        continue
                    rows = _clean_table_rows(t)
                    df = _rows_to_df(rows)
                    if df is None or df.empty:
                        continue
                    tcount += 1
                    out.append({"name": f"Tabla_PDF_P{pi:03d}_{tj:02d}", "df": df})
                    if tcount >= max_tables:
                        return out
    except Exception:
        return []
    return out


# =============================
# Image extraction (MEJORADO)
# =============================
def _guess_ext_from_content_type(content_type: str) -> str:
    ct = (content_type or "").lower()
    if "jpeg" in ct or "jpg" in ct:
        return "jpg"
    if "png" in ct:
        return "png"
    if "gif" in ct:
        return "gif"
    if "webp" in ct:
        return "webp"
    if "tiff" in ct or "tif" in ct:
        return "tif"
    return "png"


def extract_images_docx(path: str, outdir: str, min_bytes: int = 6_000) -> List[str]:
    try:
        from docx import Document
    except Exception:
        return []
    os.makedirs(outdir, exist_ok=True)

    out = []
    seen_hashes = set()
    idx = 1
    try:
        doc = Document(path)
        for rel in doc.part.rels.values():
            try:
                if "image" not in getattr(rel, "target_ref", ""):
                    continue
                img_bytes = rel.target_part.blob
                if not img_bytes or len(img_bytes) < min_bytes:
                    continue
                h = hash(img_bytes[:2048])
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)

                content_type = getattr(rel.target_part, "content_type", "")
                ext = _guess_ext_from_content_type(content_type)
                fname = f"img_{idx:03d}.{ext}"
                fpath = os.path.join(outdir, fname)
                with open(fpath, "wb") as f:
                    f.write(img_bytes)
                out.append(fpath)
                idx += 1
            except Exception:
                continue
    except Exception:
        return []
    return out


def extract_images_pptx(path: str, outdir: str, min_bytes: int = 6_000) -> List[str]:
    try:
        from pptx import Presentation
    except Exception:
        return []
    os.makedirs(outdir, exist_ok=True)

    out = []
    seen_hashes = set()
    idx = 1
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if not hasattr(shape, "image"):
                        continue
                    img = shape.image
                    blob = img.blob
                    if not blob or len(blob) < min_bytes:
                        continue
                    h = hash(blob[:2048])
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)

                    ext = img.ext if getattr(img, "ext", None) else "png"
                    fname = f"img_{idx:03d}.{ext}"
                    fpath = os.path.join(outdir, fname)
                    with open(fpath, "wb") as f:
                        f.write(blob)
                    out.append(fpath)
                    idx += 1
                except Exception:
                    continue
    except Exception:
        return []
    return out


def extract_images_pdf(path: str, outdir: str, min_bytes: int = 8_000) -> List[str]:
    """
    Extrae imágenes embebidas con PyMuPDF (fitz). Filtra mini-íconos.
    """
    try:
        import fitz  # PyMuPDF
    except Exception:
        return []
    os.makedirs(outdir, exist_ok=True)

    out = []
    seen_hashes = set()
    idx = 1
    try:
        doc = fitz.open(path)
        for page_i in range(len(doc)):
            page = doc[page_i]
            imgs = page.get_images(full=True)
            for img in imgs:
                try:
                    xref = img[0]
                    base = doc.extract_image(xref)
                    img_bytes = base.get("image")
                    ext = base.get("ext", "png")
                    if not img_bytes or len(img_bytes) < min_bytes:
                        continue
                    h = hash(img_bytes[:2048])
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)

                    fname = f"img_{idx:03d}_p{page_i+1:03d}.{ext}"
                    fpath = os.path.join(outdir, fname)
                    with open(fpath, "wb") as f:
                        f.write(img_bytes)
                    out.append(fpath)
                    idx += 1
                except Exception:
                    continue
        doc.close()
    except Exception:
        return []
    return out


# =============================
# Draft creation (MEJORADO)
# =============================
def convert_file_to_draft(uploaded_file):
    # guarda el archivo subido a disco temporal
    with tempfile.NamedTemporaryFile(delete=False, suffix=Path(uploaded_file.name).suffix) as tmp:
        tmp.write(uploaded_file.read())
        tmp_path = tmp.name

    filename = uploaded_file.name
    ext = Path(filename).suffix.lower().lstrip(".")

    tmp_dir = Path(tempfile.mkdtemp())
    images_dir = tmp_dir / "images"
    attachments_dir = tmp_dir / "attachments"
    images_dir.mkdir(exist_ok=True)
    attachments_dir.mkdir(exist_ok=True)

    md_text = ""
    excel_sheets = []
    extracted_tables: List[Dict[str, Any]] = []

    # Tablas estructuradas (excel/csv)
    if ext == "csv":
        df = extract_table_from_csv(tmp_path)
        md_text = "# " + Path(filename).stem + "\n\n" + df_to_markdown_table(df) + "\n"
    elif ext in ["xls", "xlsx"]:
        xls = pd.ExcelFile(tmp_path)
        excel_sheets = xls.sheet_names
        df = extract_table_from_excel(tmp_path, excel_sheets[0] if excel_sheets else None)
        md_text = "# " + Path(filename).stem + "\n\n" + df_to_markdown_table(df) + "\n"

    # Tablas desde documentos (docx/pptx/pdf)
    if ext == "docx":
        extracted_tables = extract_tables_docx(tmp_path)
    elif ext == "pptx":
        extracted_tables = extract_tables_pptx(tmp_path)
    elif ext == "pdf":
        extracted_tables = extract_tables_pdf(tmp_path)

    # Conversión general a Markdown con MarkItDown (texto)
    try:
        converter = MarkItDown()
        result = converter.convert(tmp_path)
        md_text_from_markitdown = result.text_content or ""
        if not md_text.strip():
            md_text = md_text_from_markitdown
        else:
            # si ya había tabla excel/csv, igual guardamos el texto como "extra" por si sirve
            if md_text_from_markitdown.strip():
                md_text = md_text + "\n\n---\n\n" + md_text_from_markitdown
    except Exception:
        pass

    # Imágenes desde documentos
    if ext == "pdf":
        extract_images_pdf(tmp_path, str(images_dir))
    elif ext == "docx":
        extract_images_docx(tmp_path, str(images_dir))
    elif ext == "pptx":
        extract_images_pptx(tmp_path, str(images_dir))

    return {
        "tmp_dir": tmp_dir,
        "md_text": md_text or "",
        "original_tmp_path": tmp_path,
        "ext": ext,
        "source_filename": filename,
        "excel_sheets": excel_sheets,
        "extracted_tables": extracted_tables,  # <-- NUEVO
    }


def make_empty_draft():
    tmp_dir = Path(tempfile.mkdtemp())
    images_dir = tmp_dir / "images"
    attachments_dir = tmp_dir / "attachments"
    images_dir.mkdir(exist_ok=True)
    attachments_dir.mkdir(exist_ok=True)

    md_text = (
        "# Título de la publicación\n\n"
        "Escribe tu contenido aquí.\n\n"
        f"{CURSOR_TOKEN}\n"
    )

    return {
        "tmp_dir": tmp_dir,
        "md_text": md_text,
        "original_tmp_path": "",
        "ext": "md",
        "source_filename": "redaccion_libre.md",
        "excel_sheets": [],
        "extracted_tables": [],  # <-- NUEVO
    }


# =============================
# Editor ACE
# =============================
def ace_editor(md_key: str, height: int):
    content = st_ace(
        value=st.session_state.get(md_key, "") or "",
        language="markdown",
        theme="chrome",
        key=f"ace_{md_key}",
        height=height,
        font_size=15,
        tab_size=2,
        show_gutter=True,
        show_print_margin=False,
        wrap=True,
        auto_update=True,
    )
    if content is None:
        content = st.session_state.get(md_key, "") or ""
    st.session_state[md_key] = content
    st.session_state["preview_dirty"] = True


def toolbar(md_key: str):
    st.caption("Inserción recomendada: usa 📍 para poner `<<CURSOR>>` donde quieras insertar.")
    c = st.columns(9)
    if c[0].button("H1"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n# Título\n")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[1].button("H2"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n## Subtítulo\n")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[2].button("**B**"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "**negrita**")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[3].button("*I*"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "*cursiva*")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[4].button("Link"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "[texto](https://example.com)")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[5].button("Lista"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n- item 1\n- item 2\n")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[6].button("Código"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n```r\n# código aquí\n```\n")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[7].button("Ecuación"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n$$\nE = mc^2\n$$\n")
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c[8].button("Tabla"):
        st.session_state[md_key] = insert_at_cursor_token(st.session_state.get(md_key, ""), "\n| Col1 | Col2 |\n|---|---|\n| a | b |\n")
        st.session_state["preview_dirty"] = True
        st.rerun()

    c2 = st.columns([1, 1, 2])
    if c2[0].button("📍 Cursor"):
        st.session_state[md_key] = ensure_cursor_token(st.session_state.get(md_key, ""))
        st.session_state["preview_dirty"] = True
        st.rerun()
    if c2[1].button("🧽 Limpiar"):
        st.session_state[md_key] = remove_all_cursor_tokens(st.session_state.get(md_key, ""))
        st.session_state["preview_dirty"] = True
        st.rerun()
    c2[2].caption("Pon `<<CURSOR>>` y luego inserta imágenes/tablas donde corresponda.")


# =============================
# Crear publicación (UI)
# =============================
def crear_blog_app():
    st.title("✍️ Crear nueva publicación")

    st.session_state.setdefault("md_preview", "")
    st.session_state.setdefault("preview_auto", False)
    st.session_state.setdefault("preview_dirty", False)
    st.session_state.setdefault("cover_choice", "")

    mode = st.radio(
        "Modo de creación",
        ["Desde archivo", "Redacción libre (Markdown)"],
        horizontal=True,
        key="create_mode",
    )

    st.session_state.setdefault("editor_height", 980)
    st.slider("Tamaño del editor", 650, 1600, key="editor_height")

    if mode == "Desde archivo":
        uploaded = st.file_uploader(
            "Sube tu documento",
            type=[
                "pdf", "docx", "pptx", "txt", "md", "markdown", "rmd", "qmd", "tex", "rtf",
                "ppt", "xls", "xlsx", "csv"
            ],
            key="upload_doc",
        )
        if uploaded is None:
            st.stop()

        if st.session_state.get("last_upload_name") != uploaded.name:
            with st.spinner("Convirtiendo a Markdown + extrayendo imágenes + detectando tablas…"):
                draft = convert_file_to_draft(uploaded)

            st.session_state["draft"] = draft
            st.session_state["last_upload_name"] = uploaded.name

            base = Path(uploaded.name).stem
            st.session_state["draft_title"] = infer_title_from_md(draft["md_text"], base)
            st.session_state["draft_desc"] = infer_description_from_md(draft["md_text"])
            st.session_state["md_edit"] = draft["md_text"]
            st.session_state["md_preview"] = draft["md_text"]
            st.session_state["preview_dirty"] = False
            st.session_state["cover_choice"] = ""

        draft = st.session_state["draft"]

    else:
        if st.session_state.get("draft_free_ready") != True:
            draft = make_empty_draft()
            st.session_state["draft"] = draft
            st.session_state["draft_free_ready"] = True
            st.session_state["last_upload_name"] = "__FREE__"
            st.session_state["draft_title"] = "Nueva publicación"
            st.session_state["draft_desc"] = ""
            st.session_state["md_edit"] = draft["md_text"]
            st.session_state["md_preview"] = ""
            st.session_state["preview_dirty"] = True
            st.session_state["cover_choice"] = ""

        draft = st.session_state["draft"]
        uploaded = None

    tmp_dir: Path = draft["tmp_dir"]
    images_dir = tmp_dir / "images"
    attachments_dir = tmp_dir / "attachments"
    images_dir.mkdir(exist_ok=True)
    attachments_dir.mkdir(exist_ok=True)

    all_image_names = [p.name for p in sorted(images_dir.glob("*"))]
    extracted_tables = draft.get("extracted_tables", []) or []

    # Excel UI
    if draft.get("ext") in ["xls", "xlsx"] and draft.get("excel_sheets"):
        st.info("Excel detectado: tablas garantizadas en Markdown.")
        sheet = st.selectbox("Hoja a convertir", draft["excel_sheets"], key="excel_sheet_pick")
        try:
            df = extract_table_from_excel(draft["original_tmp_path"], sheet)
            with st.expander("👀 Vista previa tabla (Excel)", expanded=False):
                st.dataframe(df.head(200), use_container_width=True)
            if st.button("🔁 Reemplazar editor con tabla de esta hoja", use_container_width=True):
                md_table = "# " + Path(draft["source_filename"]).stem + f" — {sheet}\n\n" + df_to_markdown_table(df) + "\n"
                st.session_state["md_edit"] = md_table
                st.session_state["preview_dirty"] = True
                st.rerun()
        except Exception as e:
            st.warning(f"No pude leer la hoja: {e}")

    col_left, col_editor, col_side = st.columns([4, 6, 3], vertical_alignment="top")

    with col_left:
        st.subheader("📄 Fuente / Documento")
        with st.container(border=True):
            if mode == "Desde archivo":
                if uploaded and uploaded.name.lower().endswith(".pdf") and pdf_viewer is not None:
                    pdf_viewer(draft["original_tmp_path"])
                else:
                    st.markdown("Vista previa (texto recortado).")
                    st.write((st.session_state.get("md_edit", "") or "")[:2500])
            else:
                st.info("Estás en **Redacción libre**. No hay documento fuente.")

    with col_editor:
        st.subheader("📝 Editor Markdown")

        title = st.text_input("Título", value=st.session_state.get("draft_title", ""), key="title_mode1")
        desc = st.text_area("Descripción", value=st.session_state.get("draft_desc", ""), height=110, key="desc_mode1")

        toolbar("md_edit")

        tabs = st.tabs(["✍️ Editor", "👁️ Preview (Lectura grande)"])
        with tabs[0]:
            with st.container(border=True):
                ace_editor("md_edit", height=st.session_state["editor_height"])

        with tabs[1]:
            cA, cB, cC = st.columns([1.2, 1.2, 2], vertical_alignment="bottom")
            with cA:
                if st.button("🔄 Actualizar vista previa", use_container_width=True):
                    st.session_state["md_preview"] = st.session_state.get("md_edit", "") or ""
                    st.session_state["preview_dirty"] = False
                    st.rerun()
            with cB:
                st.toggle("Auto-preview", value=st.session_state.get("preview_auto", False), key="preview_auto")
            with cC:
                if st.session_state.get("preview_dirty", False):
                    st.warning("Vista previa desactualizada (tienes cambios sin aplicar).")
                else:
                    st.success("Vista previa actualizada.")

            if st.session_state.get("preview_auto", False):
                st.session_state["md_preview"] = st.session_state.get("md_edit", "") or ""
                st.session_state["preview_dirty"] = False

            with st.container(border=True):
                st.markdown(st.session_state.get("md_preview", "") or "", unsafe_allow_html=True)

        # =============================
        # NUEVO: Tablas detectadas en el documento
        # =============================
        if mode == "Desde archivo" and extracted_tables:
            st.markdown("##### 📋 Tablas detectadas en el documento")
            names = [t["name"] for t in extracted_tables]
            pick_t = st.selectbox("Selecciona una tabla detectada", names, key="pick_detected_table")
            t_obj = next((x for x in extracted_tables if x["name"] == pick_t), None)
            if t_obj is not None:
                df_t = t_obj.get("df")
                with st.expander("👀 Vista previa (tabla detectada)", expanded=False):
                    if isinstance(df_t, pd.DataFrame):
                        st.dataframe(df_t.head(200), use_container_width=True)
                tt = st.text_input("Título (opcional)", value=pick_t, key="det_table_title")
                if st.button("➕ Insertar esta tabla en el cursor", use_container_width=True, key="btn_insert_detected_table"):
                    snippet = make_table_snippet(df_t, title=tt or pick_t)
                    st.session_state["md_edit"] = insert_at_cursor_token(st.session_state.get("md_edit", "") or "", snippet)
                    st.session_state["preview_dirty"] = True
                    st.rerun()

        st.markdown("##### 🖼️ Insertar imagen en el cursor")

        add_imgs = st.file_uploader(
            "Subir imágenes (se guardan en images/ del draft)",
            type=["png", "jpg", "jpeg", "webp"],
            accept_multiple_files=True,
            key="imgs_inline_uploader",
        )
        if add_imgs:
            for f in add_imgs:
                try:
                    (images_dir / f.name).write_bytes(f.read())
                except Exception:
                    pass
            st.success("Imágenes agregadas.")
            st.rerun()

        all_image_names = [p.name for p in sorted(images_dir.glob("*"))]

        if all_image_names:
            sel_img = st.selectbox("Selecciona una imagen", all_image_names, key="sel_img_mode1_simple")
            pprev = images_dir / sel_img
            if pprev.exists():
                st.image(str(pprev), use_container_width=True)

            if st.button("➕ Insertar imagen en el cursor", use_container_width=True):
                st.session_state["md_edit"] = insert_at_cursor_token(
                    st.session_state.get("md_edit", "") or "",
                    make_image_snippet(sel_img)
                )
                st.session_state["preview_dirty"] = True
                st.rerun()

        st.markdown("##### 📊 Insertar tabla (CSV/Excel) en el cursor")

        table_up = st.file_uploader(
            "Sube un CSV o Excel para convertirlo a tabla Markdown e insertarlo donde está el cursor",
            type=["csv", "xlsx", "xls"],
            key="table_uploader_inline",
        )

        if table_up is not None:
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(table_up.name).suffix) as t:
                t.write(table_up.read())
                tpath = t.name

            try:
                ext2 = Path(table_up.name).suffix.lower()
                if ext2 == ".csv":
                    df = extract_table_from_csv(tpath)
                else:
                    xls = pd.ExcelFile(tpath)
                    sheet_name = st.selectbox("Hoja (Excel)", xls.sheet_names, key="inline_table_sheet_pick")
                    df = extract_table_from_excel(tpath, sheet_name)

                with st.expander("👀 Vista previa tabla", expanded=False):
                    st.dataframe(df.head(200), use_container_width=True)

                tab_title = st.text_input("Título de la tabla (opcional)", value="Tabla", key="inline_table_title")
                if st.button("➕ Insertar tabla en el cursor", use_container_width=True, key="insert_table_btn"):
                    snippet = make_table_snippet(df, title=tab_title or "Tabla")
                    st.session_state["md_edit"] = insert_at_cursor_token(st.session_state.get("md_edit", "") or "", snippet)
                    st.session_state["preview_dirty"] = True
                    st.rerun()

            except Exception as e:
                st.warning(f"No pude leer esa tabla: {e}")
            finally:
                try:
                    os.remove(tpath)
                except Exception:
                    pass

        st.markdown("##### 📎 Adjuntar archivos")
        att = st.file_uploader("Adjuntos", accept_multiple_files=True, key="att_mode1")
        if att:
            for f in att:
                (attachments_dir / f.name).write_bytes(f.read())
            st.success("Adjuntos agregados.")
            st.rerun()

        d1, d2 = st.columns([1, 1], vertical_alignment="bottom")
        md_bytes = (st.session_state.get("md_edit", "") or "").encode("utf-8")
        d1.download_button(
            "📥 Descargar MD",
            data=md_bytes,
            file_name=f"{slugify(title)}.md",
            mime="text/markdown",
            use_container_width=True
        )

        draft_pkg_dir = tmp_dir / "_draft_package"
        draft_pkg_dir.mkdir(exist_ok=True)
        (draft_pkg_dir / "post.md").write_text(st.session_state.get("md_edit", "") or "", encoding="utf-8")

        (draft_pkg_dir / "images").mkdir(exist_ok=True)
        for p in images_dir.glob("*"):
            shutil.copy2(p, draft_pkg_dir / "images" / p.name)

        any_att = list(attachments_dir.glob("*"))
        if any_att:
            (draft_pkg_dir / "attachments").mkdir(exist_ok=True)
            for p in any_att:
                shutil.copy2(p, draft_pkg_dir / "attachments" / p.name)

        d2.download_button(
            "📦 Descargar ZIP (MD+imágenes+adjuntos)",
            data=zip_folder(draft_pkg_dir),
            file_name=f"{slugify(title)}_draft.zip",
            mime="application/zip",
            use_container_width=True
        )

        st.divider()
        if st.button("📌 Publicar", type="primary", use_container_width=True):
            created_at = datetime.now().strftime("%Y-%m-%d %H:%M")
            post_id = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{slugify(title)[:40]}"
            post_dir = BLOG_DIR / post_id
            post_dir.mkdir(parents=True, exist_ok=True)

            md_final = st.session_state.get("md_edit", "") or ""
            (post_dir / "post.md").write_text(md_final, encoding="utf-8")

            (post_dir / "images").mkdir(exist_ok=True)
            for p in images_dir.glob("*"):
                shutil.copy2(p, post_dir / "images" / p.name)

            any_att = list(attachments_dir.glob("*"))
            if any_att:
                (post_dir / "attachments").mkdir(exist_ok=True)
                for p in any_att:
                    shutil.copy2(p, post_dir / "attachments" / p.name)

            cover_value = ""
            if st.session_state.get("cover_choice"):
                cover_value = f"images/{st.session_state['cover_choice']}"

            meta = {
                "id": post_id,
                "title": title.strip() or "Sin título",
                "description": (desc or "").strip(),
                "created_at": created_at,
                "cover": cover_value,
                "source_filename": draft.get("source_filename", "redaccion_libre.md"),
            }
            upsert_post_meta(meta)

            st.success("Publicado ✅")
            st.info("Ahora puedes ir a **Boletines** para ver el post.")

    with col_side:
        st.subheader("🖼️ Imágenes (ver / descargar / portada)")
        all_image_names = [p.name for p in sorted(images_dir.glob("*"))]

        with st.expander("⭐ Portada del post", expanded=True):
            if not all_image_names:
                st.info("Aún no hay imágenes. Sube una para poder elegir portada.")
                st.session_state["cover_choice"] = ""
            else:
                current = st.session_state.get("cover_choice", "")
                options = ["(Sin portada)"] + all_image_names
                idx = 0
                if current in all_image_names:
                    idx = options.index(current)

                pick = st.selectbox("Elige la imagen de portada", options, index=idx, key="cover_picker")
                if pick == "(Sin portada)":
                    st.session_state["cover_choice"] = ""
                else:
                    st.session_state["cover_choice"] = pick

                if st.session_state.get("cover_choice"):
                    cp = images_dir / st.session_state["cover_choice"]
                    if cp.exists():
                        st.image(str(cp), use_container_width=True)
                        st.caption(f"Portada seleccionada: {cp.name}")

        if all_image_names:
            with st.container(border=True):
                for name in all_image_names:
                    p = images_dir / name
                    if p.exists():
                        st.image(str(p), use_container_width=True)
                        st.download_button(
                            f"⬇️ Descargar {name}",
                            data=p.read_bytes(),
                            file_name=name,
                            use_container_width=True,
                            key=f"dl_draft_img_{name}",
                        )
                        st.divider()
        else:
            st.info("No hay imágenes aún. Súbelas abajo o en el editor.")

        st.markdown("##### ➕ Subir imágenes extra")
        imgs_up = st.file_uploader(
            "Imágenes",
            type=["png", "jpg", "jpeg", "webp"],
            accept_multiple_files=True,
            key="extra_imgs_side",
        )
        if imgs_up:
            for f in imgs_up:
                (images_dir / f.name).write_bytes(f.read())
            st.success("Imágenes agregadas.")
            st.rerun()


def main():
    crear_blog_app()


if __name__ == "__main__":
    crear_blog_app()
